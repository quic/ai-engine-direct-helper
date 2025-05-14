from langchain_core.callbacks import CallbackManager, BaseCallbackHandler
from langchain_core.prompts import PromptTemplate
from langchain_core.document_loaders import BaseLoader

from langchain.chains.summarize import load_summarize_chain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import Docx2txtLoader, PyPDFLoader, PythonLoader, TextLoader, JSONLoader
from langchain_community.document_loaders import DirectoryLoader

from langchain.schema import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

from ChainUtils import GenieLLM, GenieLLMCallbackHandler
from gradio import ChatMessage
import threading
import time
import os


###########################################################################

PROMPT_TEMPLATE = (
    "你的工作是用中文写出以下文档的摘要:\n"
    "================文档内容开始================\n"
    "{text}\n"
    "================文档内容结束================\n"
)

REFINE_TEMPLATE = (
    "你的工作是编写最终摘要\n"
    "我们已经提供了一定程度的现有摘要：\n"
    "================现有摘要开始================\n"
    "{existing_answer\n}"
    "================现有摘要结束================\n"
    "我们有机会完善现有的摘要"
    "(only if needed) 下面有更多背景信息：\n"
    "================新的背景信息开始================\n"
    "{text}\n"
    "================新的背景信息结束================\n"
    "鉴于新的背景信息，完善摘要。\n"
)

###########################################################################

# Remove blank line.
def remove_blank_line(text):
    text = text.strip()
    lines = [line for line in text.splitlines() if line.strip()]
    text = '\n'.join(lines)
    return text

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.group import GroupShape

# https://python.langchain.ac.cn/docs/how_to/document_loader_custom/
# https://blog.csdn.net/weixin_40941102/article/details/145512744
class PPTXLoader(BaseLoader):
    def __init__(self, file_path):
        self.file_path = file_path
        self.content = ""

    def extract_text_frame_info(self, text_frame):
        if text_frame and hasattr(text_frame, 'paragraphs'):
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    self.content += " " + run.text

    def extract_table_info(self, table):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                self.extract_text_frame_info(cell.text_frame)

    def extract_shape_info(self, shape):
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            self.extract_table_info(shape.table)
        elif shape.has_text_frame:
            self.extract_text_frame_info(shape.text_frame)
        elif isinstance(shape, GroupShape):
            for sub_shape in shape.shapes:
                self.extract_shape_info(sub_shape)

    def extract_pptx_info(self, pptx_path):
        prs = Presentation(pptx_path)
        for slide_index, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    self.extract_table_info(shape.table)
                elif shape.has_text_frame:
                    self.extract_text_frame_info(shape.text_frame)
                elif isinstance(shape, GroupShape):
                    for sub_shape in shape.shapes:
                        self.extract_shape_info(sub_shape)

    def load(self):
        self.extract_pptx_info(self.file_path)
        documents = []
        documents.append(Document(page_content=self.content, metadata={"source": self.file_path}))
        return documents

###########################################################################

class DocUtils():
    def __init__(self, 
        llm: GenieLLM = None) -> None:
        self.llm = llm
        print(self.llm.model_name)

class CustomChainCallbackHandler(BaseCallbackHandler, GenieLLMCallbackHandler):
    def __init__(self, chain_type, chain_length):
        self.chain_type = chain_type
        self.chain_step_count = 0
        self.chain_length = chain_length
        self.stream_chunk = ""
        self.step_lock = threading.Lock()
        # self.stream_lock = threading.Lock()

    def on_query_start(self, serialized, prompts, **kwargs):
        time.sleep(0.05)  # Waiting for the stream to be readed by the main thread.
        self.step_lock.acquire()
        self.chain_step_count += 1
        self.step_lock.release()
        self.stream_chunk = ""

        #print("[ ", self.chain_step_count, self.chain_length, " ]")

        #print("\n\n===== 语言模型GenieLLM开始处理 =====\n\n")
        print(f"语言模型开始处理，提示信息: {prompts}")
        #print(f"步骤 {self.chain_step_count} 开始，当前处理的提示信息:", prompts)
        #if self.chain_step_count == self.chain_length:
        #    print("\nThe last step begin.\n")
        #if self.chain_type == 'map_reduce':
        #    print("\n===== MapReduce begin =====\n")

    def on_query_end(self, response, **kwargs):
        #print("\n\n===== 语言模型GenieLLM结束处理 =====\n\n")
        #print(f"语言模型结束处理，提示信息: {response}")
        #print(f"M步骤 {self.chain_step_count} 结束，当前步骤结果:", response)
        if self.chain_step_count == self.chain_length:
            time.sleep(0.05)  # Waiting for the stream to be readed by the main thread.
            self.step_lock.acquire()
            self.chain_step_count += 1
            self.step_lock.release()
            #print("\nThe last step end.\n")
        #if self.chain_type == 'map_reduce':
            #print("\n===== MapReduce end =====\n")

    def on_llm_start(self, serialized, prompts, **kwargs):
        #print("\n\n===== 语言模型开始处理 =====\n\n")
        #print(f"语言模型开始处理，提示信息: {prompts}")
        pass

    def on_llm_end(self, response, **kwargs):
        #print("\n\n===== 语言模型处理结束 =====\n\n")
        #print(f"语言模型处理结束，响应结果: {response}")
        pass

    def on_llm_new_token(self, token, **kwargs):
        #print("\n\n===== new token =====\n\n")
        #print(f"新生成的 token: {token}")
        #print(token, end="", flush=True)
        self.stream_chunk += token

    def on_llm_error(self, error, **kwargs):
        print("\n\n===== 语言模型处理出错 =====\n\n")
        print(f"语言模型处理出错，错误信息: {error}")
        pass

    def on_chain_start(self, serialized, inputs, **kwargs):
        #print("\n\n===== 链开始执行 =====\n\n")
        #print(f"链开始执行，输入内容: {inputs}")
        #if 'map_reduce' in serialized['name']:
        #    print("MapReduce 链开始执行")
        pass

    def on_chain_end(self, outputs, **kwargs):
        #print("\n\n===== 链执行结束 =====\n\n")
        #print(f"链执行结束，输出结果: {outputs}")
        #if 'map_reduce' in outputs['name']:
        #    print("MapReduce 链执行结束")
        pass

class DocSummarize():
 
    def __init__(self, 
        llm: GenieLLM = None) -> None:
        self.llm = llm

    def is_chinese(self):
        print(self.llm.model_name)
        if "Qwen" in self.llm.model_name:
            return True
        return False

    def load_docs(self, file_path):

        # TODO: load multipal files. load different files format. like docx, txt, pdf, etc.
        # [".pdf", ".docx", ".pptx", ".txt", ".md", ".py", ".c", ".cpp", ".h", ".hpp" ]

        documents = []

        loader = DirectoryLoader(file_path, glob="**/*.pdf", loader_cls=PyPDFLoader, use_multithreading=True, loader_kwargs={"mode": 'single'})
        document = loader.load()
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.docx", loader_cls=Docx2txtLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.pptx", loader_cls=PPTXLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.txt", loader_cls=TextLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.md", loader_cls=TextLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.c*", loader_cls=TextLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.h*", loader_cls=TextLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        loader = DirectoryLoader(file_path, glob="**/*.py", loader_cls=PythonLoader, use_multithreading=True)
        document = loader.load()            
        documents += document

        #print(documents)

        text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=3000, chunk_overlap=200,
                        separators=["\n\n", "\n", "。", "！", "？", ".", "?", "!", " ", ""],
                        length_function=self.llm.get_num_tokens
                    )

        self.split_documents = text_splitter.split_documents(documents)
        self.chain_length = len(self.split_documents)
        #print(self.split_documents)
        print(self.chain_length)

    def init(self, prompt, chain_type="map_reduce", file_path=None):
        self.prompt = prompt

        if not os.path.exists(file_path):
            return False

        self.file_path = file_path
        self.load_docs(file_path)
        if self.chain_length == 0:
            return False

        prompt_template = PROMPT_TEMPLATE
        refine_template = REFINE_TEMPLATE

        additional_info = f"编写摘要的附加说明或要求：{prompt}\n"
        prompt_info = "中文摘要:"

        if len(prompt) > 3:
            prompt_template += additional_info
            refine_template += additional_info

        prompt_template += prompt_info
        refine_template += prompt_info

        PROMPT = PromptTemplate(template=prompt_template, input_variables=["text"])

        self.chain = None
        intermediate_steps = False

        if (chain_type == "map_reduce") and (self.chain_length > 1):
            self.chain_length += 1
        elif self.chain_length == 1:
            chain_type = "stuff"

        self.callback_handler = CustomChainCallbackHandler(chain_type, self.chain_length)
        self.llm.set_callback(self.callback_handler)

        # https://github.com/langchain-ai/langchain/blob/master/libs/langchain/langchain/chains/combine_documents/
        # https://github.com/langchain-ai/langchain/tree/master/libs/langchain/langchain/chains/summarize/
        if chain_type == "map_reduce":  # 各段文档同时处理，之后汇总。
            if self.is_chinese():
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, return_intermediate_steps=intermediate_steps,
                                            map_prompt=PROMPT, combine_prompt=PROMPT, token_max=4096)
            else:
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, return_intermediate_steps=intermediate_steps,
                                            token_max=4096)

        elif chain_type == "refine":    # 将前一段处理后得到的摘要作为后一段的输入，依次处理直到各段文档全部完成。
            if self.is_chinese():
                REFINE_PROMPT = PromptTemplate.from_template(template=refine_template)
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, return_intermediate_steps=intermediate_steps,
                                                  question_prompt=PROMPT, refine_prompt=REFINE_PROMPT, token_max=4096)
            else:
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, return_intermediate_steps=intermediate_steps,
                                                  token_max=4096)
        else: # stuff
            self.chain_length = 1
            if self.is_chinese():
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, # return_intermediate_steps=intermediate_steps,
                                                  prompt=PROMPT)
            else:
                self.chain = load_summarize_chain(self.llm, chain_type=chain_type, verbose=False, # return_intermediate_steps=intermediate_steps,
                                                  prompt=PROMPT)

        return True

    def chain_thread(self):
        output = self.chain.invoke({"input_documents": self.split_documents, 'token_max': 4096}, {"callbacks": [self.callback_handler]}, return_only_outputs=True)
        self.summarize_text = output['output_text']

    def stop(self):
        self.llm.stop()

    def summarize(self, chatbot=None, max_length=2048, temp=0.8, top_k=40, top_p=0.95):
        self.summarize_thread = threading.Thread(target=self.chain_thread)
        self.summarize_thread.start()

        chain_step_count = 0
        stream_text = ""
        step_change = False

        while chain_step_count <= self.chain_length:
            self.callback_handler.step_lock.acquire()

            if chain_step_count < self.callback_handler.chain_step_count:
                chain_step_count = self.callback_handler.chain_step_count
                step_change = True

                if chain_step_count > self.chain_length:
                    self.callback_handler.step_lock.release()
                    break

            self.callback_handler.step_lock.release()

            if step_change and chain_step_count > 1:
                time_to_first_token, prompt_speed, eval_speed, profile = self.llm.get_profile_str()
                # print(profile)
                chatbot[-1].content += "<br>" + profile
                yield chatbot, "", "", ""

            if step_change:
                print("<Setp>", chain_step_count, " / ", self.chain_length)  # Step changed.
                if chain_step_count == self.chain_length: # 文档最后的总结。
                    chatbot.append(ChatMessage(role="assistant", content=""))
                    pass
                else:   # 文档总结的中间阶段。
                    title = "<big><b>{ 第 " + str(chain_step_count) + " 部分总结 | 共 " + str(self.chain_length - 1) + " 个部分 }</b></big>"
                    chatbot.append(ChatMessage(role="assistant", content="", metadata={"title": title}))
                yield chatbot, "", "", ""

            step_change = False

            if chain_step_count > 0:  # generating ...
                stream_text = self.callback_handler.stream_chunk
                if len(stream_text) > 0:
                    #print("<Thinking>", chain_step_count)
                    #print(stream_text)
                    chatbot[-1].content = stream_text
                    yield chatbot, "", "", ""

            time.sleep(0.01)

        time_to_first_token, prompt_speed, eval_speed, profile = self.llm.get_profile_str()
        # print(profile)
        try:
            chatbot[-1].content += "<br>" + profile
        except Exception as e:
            print(e)

        self.summarize_thread.join()

        summarize_text = self.summarize_text
        # print(summarize_text)
        yield chatbot, time_to_first_token, prompt_speed, eval_speed
        #return summarize_text


def main():
    import argparse

    parser = argparse.ArgumentParser()
    parser.add_argument("--stream", action="store_true")
    parser.add_argument("--prompt", default="你好", type=str)
    args = parser.parse_args()
    prompt = args.prompt

    model_name = "Qwen2.0-7B-AR32-v29"

    llm = GenieLLM()
    llm.init(model_name=model_name)
    llm.set_params(str(4096), str(1), str(1), str(1))
    #llm.set_params(str(4096), str(0.8), str(40), str(0.95))

    chain_type = "map_reduce"
    doc_summarize = DocSummarize(llm)
    doc_summarize.init(prompt=prompt, chain_type=chain_type, file_path="files")
    text = doc_summarize.summarize()
    return
    #print(text)

    round = 1
    for i in range(round):
        print("\n\n")
        print("=" * 60)
        print("\n\n")

        llm.set_params(str(4096), str(1), str(1), str(1))
        doc_summarize = DocSummarize(llm)
        doc_summarize.init(prompt=prompt, chain_type=chain_type, file_path="files")
        text = doc_summarize.summarize()
        #print(text)

if __name__ == "__main__":

    main()
